from pathlib import Path
import shutil
import textwrap

from PIL import Image
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(r"D:\projects\maopaoweic-site")
OUT = ROOT / "outputs" / "kindergarten_observation"
ASSETS = OUT / "assets"

SOURCE_IMAGES = [
    Path(r"D:\软件\微信\微信聊天储存位置\xwechat_files\wxid_pv444vjqjy7522_2e41\temp\RWTemp\2026-05\255dc5dc6bb150ee25d5f70427d60cd9\2ae6c7a923ae01dc86ff6defeaaf00f6.jpg"),
    Path(r"D:\软件\微信\微信聊天储存位置\xwechat_files\wxid_pv444vjqjy7522_2e41\temp\RWTemp\2026-05\255dc5dc6bb150ee25d5f70427d60cd9\fc9e41976dd399e821b078df42d898d1.jpg"),
    Path(r"D:\软件\微信\微信聊天储存位置\xwechat_files\wxid_pv444vjqjy7522_2e41\temp\RWTemp\2026-05\255dc5dc6bb150ee25d5f70427d60cd9\a4cd357f68812f8c0ed5c12edd65440e.jpg"),
    Path(r"D:\软件\微信\微信聊天储存位置\xwechat_files\wxid_pv444vjqjy7522_2e41\temp\RWTemp\2026-05\255dc5dc6bb150ee25d5f70427d60cd9\6a55bccb452931c29154ed6b4898b287.jpg"),
    Path(r"D:\软件\微信\微信聊天储存位置\xwechat_files\wxid_pv444vjqjy7522_2e41\temp\RWTemp\2026-05\255dc5dc6bb150ee25d5f70427d60cd9\7b7471d6eea95eb3e3f8d50797aba03b.jpg"),
]


TITLE = "小班后期幼儿大型建构游戏连续观察案例"
SUBTITLE = "从简单摆弄到情境游戏的发生"


STAGES = [
    {
        "name": "第一阶段：材料在场，游戏还没有真正发生",
        "period": "区域开放初期",
        "behavior": "户外区域开放后，幼儿会自然走进大型建构材料旁边。几个孩子围着灰色框架和橙色、薄荷绿面板看一看、摸一摸，有的把三角板搬起来又放下，有的从材料旁边走过去，很快转向旁边的钻爬架、球和小车。大型积木虽然摆在场地中央，但真正停留下来搭建的幼儿不多。个别幼儿会把一块面板拖到身边，反复翻面、拍打中间的圆孔，更多是在熟悉材料的重量、形状和连接方式。",
        "teacher": "教师没有急着要求幼儿“搭一个什么”，先观察孩子愿意停留在哪里、哪些材料容易被拿起。随后把大型积木从分散堆放调整为低矮、可见、便于取放的状态，同时增加蓝色坐垫、软球、小推车等辅助材料。教师和幼儿一起商量简单规则：大的板块两个人搬，搭好的地方可以进去玩，想拆的时候要先问一问正在玩的同伴。",
        "analysis": "这一阶段幼儿的兴趣更多来自材料本身，游戏目的还不清楚。大型材料体积大、连接方式不熟悉，对小班后期幼儿来说既有吸引力，也有一点门槛。教师的支持重点不是讲解作品，而是降低进入难度，让材料看起来“可以玩、敢去玩”。规则协商也让幼儿知道这里不仅能搬，还能一起建、一起玩。",
        "focus": ["游戏技能：以感知和摆弄为主", "合作能力：并行游戏多，协商少", "语言交流：多为“我要这个”“给我”", "游戏情境：尚未形成稳定主题"],
    },
    {
        "name": "第二阶段：开始搭建，但方式比较单一",
        "period": "自由探索阶段",
        "behavior": "材料调整后，幼儿停留时间明显变长。孩子们开始把长方形面板一块接一块平铺在地上，或者用橙色方板围出一个低矮的圈。有的幼儿把三角板斜靠在方板旁边，说“这个放这里”，但很快发现会倒，又重新换位置。两个孩子会同时搬一块大板，一个说“你拿那边”，另一个跟着移动。搭建作品多数是围栏、矮墙、长长的一排，连接和转角还不稳定。",
        "teacher": "教师以观察支持为主，只在孩子遇到明显困难时靠近，例如提醒“这块有点大，可以请朋友一起抬吗？”“如果想围起来，缺口可以留在哪里？”教师把相同形状的面板放在同一区域，把三角形、窗框板、坐垫放在旁边，便于幼儿看见并选择。",
        "analysis": "幼儿已经从单纯摆弄进入初步建构，能够尝试平铺、围合、简单组合。合作开始出现，但多围绕搬运和占有材料展开，还没有形成共同计划。幼儿会看同伴怎么放，再照着放一块，说明同伴经验开始成为学习资源。此时材料分类和取放方式对游戏影响明显，材料清楚，孩子就更容易延续自己的想法。",
        "focus": ["游戏技能：平铺、围合、简单连接", "合作能力：出现共同搬运", "语言交流：开始有位置和动作语言", "游戏情境：从“摆材料”走向“像个地方”"],
    },
    {
        "name": "第三阶段：教师进入游戏，帮助作品站起来",
        "period": "共同搭建阶段",
        "behavior": "当幼儿想把围合的墙面搭高时，常常出现连接不牢、屋顶放不上去、门洞被堵住的问题。图片中可以看到，教师站在一旁扶住大块屋顶板，几名幼儿围在房子旁边，有的蹲下看连接处，有的把三角板递过来。孩子们开始说“这里是门”“这个是窗户”，但搭建方式仍比较统一，常常模仿教师刚刚摆放的样子：一面墙、一块窗框、再加屋顶。",
        "teacher": "教师通过提问和示范介入：“你们想让谁进去？”“门留在哪里比较方便？”“屋顶太重，谁来帮我扶一下？”示范时不一次性搭完整，而是慢慢让幼儿参与关键步骤，如扶板、递材料、试着对齐接口。对容易倒的结构，教师用身体和手短暂支撑，等孩子找到连接位置后再松开。",
        "analysis": "教师介入让幼儿看见了更稳定的建构方式，也让他们体验到“房子可以真的钻进去”。不过观察中也发现，幼儿容易把教师示范当成唯一答案，作品外形变得相似，辅助材料使用率不高。角色和情境还停留在口头命名，真正的持续互动不多。下一步需要从“帮孩子搭成”转向“让孩子自己决定怎么用”。",
        "focus": ["游戏技能：学习立面、屋顶、门窗组合", "合作能力：围绕共同作品聚集", "语言交流：出现命名和简单分工", "游戏情境：房子、门、窗开始被赋予意义"],
    },
    {
        "name": "第四阶段：教师放手，情境游戏自然长出来",
        "period": "自主拓展阶段",
        "behavior": "有了前期经验后，幼儿不再只等教师来搭。几个孩子会先把大板围起来，再商量“这里可以进去”“这个当车门”。图片中，大型积木被搭成像小房子、像车厢的空间，孩子们钻进去坐下，有的从窗口探头，有的拿着坐垫放在入口处当“上车的地方”。两个女孩坐在围合空间里互相整理头发，像在家里照顾朋友；另一组孩子坐进“车厢”，说“开汽车啦”“我们出发”“我要坐这里”。有人想进来时，里面的孩子会说“从门这里进”。",
        "teacher": "教师逐渐退到旁边，更多用拍照、记录和简短回应支持游戏。当幼儿邀请教师时，教师以游戏角色回应，例如问“车开到哪里？”“房子里还缺什么？”而不是直接指挥搭建。材料上保留大型面板，同时持续投放坐垫、球、小车和可移动小件，让幼儿把建构作品和生活经验连接起来。",
        "analysis": "这一阶段建构目的明显清楚，幼儿不是为了把材料摆高，而是为了得到一个能坐、能躲、能进出的空间。合作从共同搬运发展到共同使用和共同维护，语言也从动作指令扩展到角色对话。材料调整的价值在这里显现出来：大型积木提供空间框架，辅助材料提供情节线索，教师放手则给了幼儿把经验变成游戏的时间。",
        "focus": ["游戏技能：能围合出可进入空间", "合作能力：协商入口、座位和角色", "语言交流：出现连续情境语言", "游戏情境：汽车、房子、出发、照顾等生活主题"],
    },
]


REFLECTION = [
    "这次连续观察提醒我，小班后期幼儿的建构游戏不是一投放材料就会自然丰富起来。开始时孩子看似“不感兴趣”，其实很多时候是还没有找到进入材料的方式。材料太大、连接太难、规则不清楚，都会让孩子只停留在摸一摸、搬一搬。",
    "教师介入要跟着游戏状态变化。前期需要整理材料、降低难度；中期需要观察等待，只在孩子卡住时给一点支点；当孩子想把作品搭高、搭稳时，教师可以进入共同游戏，提供提问和示范；当孩子已经能把作品用起来，就要及时退后，把情节、角色和协商还给幼儿。",
    "从图片里的现场可以看到，最有价值的变化不是房子搭得多漂亮，而是孩子愿意进去、愿意等同伴、愿意说“从门这里进”“我们出发”。这说明大型建构材料真正变成了幼儿共同游戏的场地。",
]


SUMMARY = [
    "游戏技能：从触摸、搬运、平铺，逐步发展到围合、立面组合和可进入空间的搭建。",
    "合作能力：从各玩各的，发展到共同搬运、共同扶板、共同维护入口和座位。",
    "语言交流：从简单索取语言，发展到位置说明、分工语言和角色情境语言。",
    "游戏情境：从没有明确主题，发展出房子、汽车、坐车出发、躲进房子、照顾朋友等生活化情节。",
    "教师支持：从材料调整，到观察支持，再到介入指导，最后逐渐放手，支持幼儿把建构经验转化为自主游戏经验。",
]


def ensure_dirs():
    ASSETS.mkdir(parents=True, exist_ok=True)


def crop_main_photo(src: Path, dst: Path):
    img = Image.open(src).convert("RGB")
    w, h = img.size
    rows = []
    pix = img.load()
    threshold_count = int(w * 0.55)
    for y in range(h):
        count = 0
        step = max(1, w // 300)
        for x in range(0, w, step):
            r, g, b = pix[x, y]
            if max(r, g, b) > 45:
                count += step
        rows.append(count > threshold_count)

    segments = []
    start = None
    for i, val in enumerate(rows):
        if val and start is None:
            start = i
        elif not val and start is not None:
            if i - start > 120:
                segments.append((start, i))
            start = None
    if start is not None and h - start > 120:
        segments.append((start, h))

    if segments:
        y1, y2 = max(segments, key=lambda s: s[1] - s[0])
        pad = 6
        y1 = max(0, y1 - pad)
        y2 = min(h, y2 + pad)
        crop = img.crop((0, y1, w, y2))
    else:
        crop = img
    crop.save(dst, quality=94)


def prepare_images():
    cropped = []
    for idx, src in enumerate(SOURCE_IMAGES, start=1):
        raw = ASSETS / f"raw_{idx}.jpg"
        clean = ASSETS / f"photo_{idx}.jpg"
        shutil.copy2(src, raw)
        crop_main_photo(src, clean)
        cropped.append(clean)
    return cropped


def wrap(text, width=36):
    return "\n".join(textwrap.wrap(text, width=width, break_long_words=False, replace_whitespace=False))


def markdown_text():
    lines = [
        f"# {TITLE}",
        "",
        f"## {SUBTITLE}",
        "",
        "### 一、观察背景",
        "",
        "观察对象为小班后期幼儿，观察场地为户外大型建构区。区域内主要材料为可连接的大型中空积木，包括灰色框架、橙色面板、薄荷绿窗框板、蓝色三角板，以及坐垫、软球、小推车等辅助材料。观察围绕同一场地、同一种大型建构材料连续展开，重点记录幼儿从简单摆弄材料，到自主建构，再到出现角色游戏和生活情境游戏的过程。",
        "",
        "从现场照片看，材料体量较大，能够围合出幼儿可以进入、坐下、探头和躲藏的空间。孩子们在搭建中不仅要处理“怎么放稳”的问题，也会自然遇到“谁来搬、从哪里进、坐在哪里、这个地方像什么”的问题。因此，这组材料很适合观察幼儿建构技能、合作行为、语言交流和游戏情境的发展。",
        "",
        "### 二、连续观察记录",
        "",
    ]
    for stage in STAGES:
        lines.extend([
            f"## {stage['name']}",
            "",
            f"观察时段：{stage['period']}",
            "",
            "| 维度 | 记录 |",
            "| --- | --- |",
            f"| 幼儿行为表现 | {stage['behavior']} |",
            f"| 教师支持 | {stage['teacher']} |",
            f"| 观察分析 | {stage['analysis']} |",
            "",
            "发展要点：",
        ])
        lines.extend([f"- {item}" for item in stage["focus"]])
        lines.append("")

    lines.extend([
        "## 三、教师支持策略变化",
        "",
        "| 阶段 | 支持重点 | 教师角色变化 |",
        "| --- | --- | --- |",
        "| 第一阶段 | 调整材料、增加辅助材料、协商基本规则 | 从管理者转向环境准备者 |",
        "| 第二阶段 | 观察幼儿尝试，提供搬运和围合的小提示 | 做安静的观察者和轻度支持者 |",
        "| 第三阶段 | 通过提问、示范、共同搭建解决结构困难 | 进入游戏，成为共同建构者 |",
        "| 第四阶段 | 减少直接指导，用回应和记录延续情境 | 逐渐放手，成为游戏的陪伴者 |",
        "",
        "## 四、教师反思",
        "",
    ])
    lines.extend([f"{i + 1}. {text}" for i, text in enumerate(REFLECTION)])
    lines.extend([
        "",
        "## 五、幼儿发展变化总结",
        "",
    ])
    lines.extend([f"- {item}" for item in SUMMARY])
    lines.append("")
    return "\n".join(lines)


def create_docx(path: Path, images):
    doc = Document()
    styles = doc.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"].font.size = Pt(10.5)

    title = doc.add_heading(TITLE, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph(SUBTITLE)
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in sub.runs:
        run.font.color.rgb = RGBColor(91, 132, 130)

    doc.add_heading("一、观察背景", level=1)
    doc.add_paragraph(
        "观察对象为小班后期幼儿，观察场地为户外大型建构区。区域内主要材料为可连接的大型中空积木，包括灰色框架、橙色面板、薄荷绿窗框板、蓝色三角板，以及坐垫、软球、小推车等辅助材料。观察围绕同一场地、同一种大型建构材料连续展开，重点记录幼儿从简单摆弄材料，到自主建构，再到出现角色游戏和生活情境游戏的过程。"
    )
    doc.add_paragraph(
        "从现场照片看，材料体量较大，能够围合出幼儿可以进入、坐下、探头和躲藏的空间。孩子们在搭建中不仅要处理“怎么放稳”的问题，也会自然遇到“谁来搬、从哪里进、坐在哪里、这个地方像什么”的问题。"
    )
    if images:
        doc.add_picture(str(images[0]), width=Inches(5.8))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("二、连续观察记录", level=1)
    for idx, stage in enumerate(STAGES):
        doc.add_heading(stage["name"], level=2)
        doc.add_paragraph(f"观察时段：{stage['period']}")
        table = doc.add_table(rows=1, cols=2)
        table.style = "Table Grid"
        table.rows[0].cells[0].text = "维度"
        table.rows[0].cells[1].text = "记录"
        for label, key in [("幼儿行为表现", "behavior"), ("教师支持", "teacher"), ("观察分析", "analysis")]:
            row = table.add_row().cells
            row[0].text = label
            row[1].text = stage[key]
        doc.add_paragraph("发展要点：")
        for item in stage["focus"]:
            doc.add_paragraph(item, style="List Bullet")
        if idx + 1 < len(images):
            doc.add_picture(str(images[idx + 1]), width=Inches(5.8))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("三、教师支持策略变化", level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    for i, label in enumerate(["阶段", "支持重点", "教师角色变化"]):
        table.rows[0].cells[i].text = label
    rows = [
        ("第一阶段", "调整材料、增加辅助材料、协商基本规则", "从管理者转向环境准备者"),
        ("第二阶段", "观察幼儿尝试，提供搬运和围合的小提示", "做安静的观察者和轻度支持者"),
        ("第三阶段", "通过提问、示范、共同搭建解决结构困难", "进入游戏，成为共同建构者"),
        ("第四阶段", "减少直接指导，用回应和记录延续情境", "逐渐放手，成为游戏的陪伴者"),
    ]
    for row_values in rows:
        row = table.add_row().cells
        for i, value in enumerate(row_values):
            row[i].text = value

    doc.add_heading("四、教师反思", level=1)
    for text in REFLECTION:
        doc.add_paragraph(text)
    doc.add_heading("五、幼儿发展变化总结", level=1)
    for item in SUMMARY:
        doc.add_paragraph(item, style="List Bullet")
    doc.save(path)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=(75, 70, 61), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = PptInches(0.08)
    frame.margin_right = PptInches(0.08)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = PPTColor(*color)
    return box


def add_card(slide, x, y, w, h, fill=(255, 255, 255), line=(239, 224, 202), radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        PptInches(x), PptInches(y), PptInches(w), PptInches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PPTColor(*fill)
    shape.line.color.rgb = PPTColor(*line)
    shape.line.width = PptPt(1.2)
    return shape


def add_photo(slide, image_path, x, y, w, h):
    pic = slide.shapes.add_picture(str(image_path), PptInches(x), PptInches(y), width=PptInches(w))
    if pic.height > PptInches(h):
        pic.height = PptInches(h)
    pic.left = PptInches(x)
    pic.top = PptInches(y)
    return pic


def add_bg(slide, color=(255, 249, 239)):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPTColor(*color)
    for x, y, s, c in [(11.7, 0.35, 0.55, (141, 196, 190)), (0.55, 6.55, 0.5, (246, 183, 105)), (12.35, 6.65, 0.35, (116, 154, 218))]:
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(x), PptInches(y), PptInches(s), PptInches(s))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PPTColor(*c)
        shp.line.fill.background()


def add_title(slide, title, tag=None):
    if tag:
        add_textbox(slide, 0.72, 0.28, 2.2, 0.28, tag, 10, True, (255, 255, 255), PP_ALIGN.CENTER)
        pill = slide.shapes[-1]
        pill.fill.solid()
        pill.fill.fore_color.rgb = PPTColor(123, 177, 172)
        pill.line.fill.background()
    add_textbox(slide, 0.72, 0.62, 8.4, 0.55, title, 25, True, (61, 72, 71))


def create_pptx(path: Path, images):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, (255, 248, 238))
    add_card(slide, 0.65, 0.68, 5.0, 5.55, (255, 255, 255), (240, 221, 194))
    add_textbox(slide, 0.95, 1.1, 4.4, 1.5, TITLE, 30, True, (57, 69, 67))
    add_textbox(slide, 0.98, 2.76, 4.1, 0.5, SUBTITLE, 17, False, (98, 132, 128))
    add_textbox(slide, 1.02, 4.8, 3.9, 0.5, "小班后期 · 户外大型建构区 · 连续观察", 14, False, (126, 103, 80))
    add_photo(slide, images[0], 6.05, 0.65, 6.55, 5.5)
    add_textbox(slide, 6.22, 6.28, 5.7, 0.4, "大型中空积木从材料探索，逐渐发展为房子、汽车和生活化情境。", 13, False, (90, 91, 83))

    # Background
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "观察背景", "BACKGROUND")
    add_card(slide, 0.75, 1.5, 5.2, 4.7, (255, 255, 255), (236, 219, 190))
    add_textbox(slide, 1.05, 1.85, 4.55, 3.8, "观察对象为小班后期幼儿，场地为户外大型建构区。材料以灰色框架、橙色面板、薄荷绿窗框板、蓝色三角板为主，并逐步加入坐垫、软球、小车等辅助材料。\n\n观察重点放在同一场地、同一种大型材料中，幼儿游戏能力、合作行为、语言交流和游戏情境的连续变化。", 16, False, (73, 78, 72))
    add_photo(slide, images[4], 6.45, 1.45, 5.95, 4.65)

    # Path
    slide = prs.slides.add_slide(blank)
    add_bg(slide, (250, 253, 248))
    add_title(slide, "观察主线", "GAME PATH")
    steps = [("材料调整", "先让材料好拿、好看见、敢去玩"), ("观察支持", "看见幼儿的真实兴趣和困难"), ("介入指导", "用提问、示范、共同搭建搭支架"), ("逐渐放手", "把情节、角色和协商还给幼儿")]
    for i, (head, body) in enumerate(steps):
        x = 0.85 + i * 3.08
        add_card(slide, x, 2.0, 2.55, 2.55, [(255, 255, 255), (244, 252, 250), (255, 248, 239), (247, 249, 255)][i], (231, 219, 200))
        add_textbox(slide, x + 0.25, 2.35, 2.05, 0.45, head, 19, True, (68, 86, 83), PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.25, 3.05, 2.05, 0.95, body, 13, False, (93, 89, 79), PP_ALIGN.CENTER)
    add_textbox(slide, 1.25, 5.6, 10.8, 0.6, "游戏从“摆弄材料”慢慢走向“共同建构”，最后生长出可持续的角色和生活情境。", 18, True, (82, 115, 111), PP_ALIGN.CENTER)

    # Stage slides
    for i, stage in enumerate(STAGES):
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, stage["name"], f"STAGE {i + 1}")
        add_photo(slide, images[min(i, len(images) - 1)], 0.78, 1.42, 5.0, 4.0)
        add_card(slide, 6.12, 1.38, 6.45, 1.42, (255, 255, 255), (238, 221, 198))
        add_textbox(slide, 6.38, 1.58, 5.95, 0.95, wrap(stage["behavior"], 44), 12, False, (70, 74, 70))
        add_card(slide, 6.12, 3.03, 3.08, 2.2, (245, 252, 250), (210, 231, 228))
        add_textbox(slide, 6.36, 3.22, 2.62, 0.35, "教师支持", 16, True, (62, 118, 113))
        add_textbox(slide, 6.36, 3.68, 2.58, 1.15, wrap(stage["teacher"], 21), 10.5, False, (70, 74, 70))
        add_card(slide, 9.48, 3.03, 3.08, 2.2, (255, 249, 241), (240, 224, 200))
        add_textbox(slide, 9.72, 3.22, 2.62, 0.35, "观察分析", 16, True, (170, 111, 58))
        add_textbox(slide, 9.72, 3.68, 2.58, 1.15, wrap(stage["analysis"], 21), 10.5, False, (70, 74, 70))
        add_card(slide, 0.78, 5.72, 11.78, 0.78, (255, 255, 255), (232, 221, 203))
        add_textbox(slide, 1.05, 5.95, 11.1, 0.32, " · ".join(stage["focus"]), 11.5, False, (86, 96, 93), PP_ALIGN.CENTER)

    # Photo evidence
    slide = prs.slides.add_slide(blank)
    add_bg(slide, (255, 248, 238))
    add_title(slide, "图片中的游戏现场", "PHOTO NOTES")
    for idx, img in enumerate(images[:4]):
        x = 0.75 + (idx % 2) * 6.2
        y = 1.35 + (idx // 2) * 2.6
        add_card(slide, x, y, 5.65, 2.18, (255, 255, 255), (236, 219, 190))
        add_photo(slide, img, x + 0.15, y + 0.15, 2.55, 1.83)
        captions = [
            "材料围合成可进入空间，孩子开始坐进去、探头、等待同伴。",
            "围合空间里出现照顾朋友、整理头发等生活化互动。",
            "教师扶住屋顶，幼儿围在连接处观察、递材料、尝试参与。",
            "孩子开始自主处理三角板和车厢式结构，建构目的更清楚。",
        ]
        add_textbox(slide, x + 2.9, y + 0.38, 2.35, 1.25, captions[idx], 12.5, False, (78, 79, 72))

    # Strategy
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "教师支持策略的变化", "SUPPORT")
    strategy = [
        ("材料调整", "把材料分类、降低取放难度；加入坐垫、球、小车等情节线索。"),
        ("观察支持", "先看孩子怎么拿、怎么放、在哪里停留，抓住真实问题再回应。"),
        ("介入指导", "通过问题和共同搭建，让幼儿看见门、窗、屋顶、入口的关系。"),
        ("逐渐放手", "当幼儿能自主使用空间时，教师退到旁边记录和回应。"),
    ]
    for i, (head, body) in enumerate(strategy):
        add_card(slide, 0.92, 1.45 + i * 1.25, 11.5, 0.92, (255, 255, 255), (235, 221, 202))
        add_textbox(slide, 1.25, 1.65 + i * 1.25, 2.1, 0.32, head, 16, True, (73, 119, 114))
        add_textbox(slide, 3.45, 1.65 + i * 1.25, 8.4, 0.32, body, 14, False, (76, 76, 70))

    # Development analysis
    slide = prs.slides.add_slide(blank)
    add_bg(slide, (250, 253, 248))
    add_title(slide, "幼儿发展变化", "DEVELOPMENT")
    for i, item in enumerate(SUMMARY[:4]):
        x = 0.9 + (i % 2) * 6.0
        y = 1.52 + (i // 2) * 2.0
        add_card(slide, x, y, 5.4, 1.42, [(255, 255, 255), (245, 252, 250), (255, 248, 239), (247, 249, 255)][i], (232, 221, 202))
        head, body = item.split("：", 1)
        add_textbox(slide, x + 0.28, y + 0.25, 4.8, 0.35, head, 17, True, (68, 95, 91))
        add_textbox(slide, x + 0.28, y + 0.74, 4.75, 0.36, body, 12.5, False, (78, 80, 74))
    add_textbox(slide, 1.2, 6.05, 10.9, 0.45, "关键变化：大型建构材料从“物”变成了幼儿共同生活和想象的场地。", 17, True, (155, 105, 63), PP_ALIGN.CENTER)

    # Reflection
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "教师反思", "REFLECTION")
    for i, text in enumerate(REFLECTION):
        add_card(slide, 0.85, 1.45 + i * 1.55, 11.65, 1.08, (255, 255, 255), (235, 221, 202))
        add_textbox(slide, 1.15, 1.66 + i * 1.55, 11.0, 0.58, text, 12.5, False, (76, 76, 70))

    # Ending
    slide = prs.slides.add_slide(blank)
    add_bg(slide, (255, 248, 238))
    add_photo(slide, images[0], 0.75, 0.75, 5.9, 4.85)
    add_textbox(slide, 7.0, 1.25, 5.1, 1.0, "看见游戏慢慢长出来", 30, True, (61, 72, 71))
    add_textbox(slide, 7.05, 2.75, 4.95, 1.25, "当材料更容易进入，教师支持更贴近现场，幼儿就会把建构经验、同伴关系和生活想象连在一起。", 17, False, (85, 91, 84))
    add_card(slide, 7.05, 4.75, 4.75, 0.78, (245, 252, 250), (213, 232, 227))
    add_textbox(slide, 7.34, 4.99, 4.15, 0.26, "材料 · 经验 · 合作 · 情境", 16, True, (74, 120, 114), PP_ALIGN.CENTER)

    prs.save(path)


def main():
    ensure_dirs()
    images = prepare_images()
    md_path = OUT / "小班大型建构游戏连续观察记录.md"
    docx_path = OUT / "小班大型建构游戏连续观察记录.docx"
    pptx_path = OUT / "小班大型建构游戏教研展示.pptx"
    md_path.write_text(markdown_text(), encoding="utf-8")
    create_docx(docx_path, images)
    create_pptx(pptx_path, images)
    print(md_path)
    print(docx_path)
    print(pptx_path)


if __name__ == "__main__":
    main()
